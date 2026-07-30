# -*- coding: utf-8 -*-
"""
三件套 PPTX 生成模板 — 深色科技风 16:9
用法：复制此文件 → 替换 slide 内容 → 改 out 路径 → 运行
环境：C:\Users\Administrator\.workbuddy\binaries\python\envs\default\Scripts\python.exe
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ======== 配色常量 ========
BG     = RGBColor(0x07,0x0b,0x16)   # 深空蓝底
CARD   = RGBColor(0x11,0x1a,0x2e)   # 卡片底色
CARD2  = RGBColor(0x16,0x21,0x3a)   # 卡片二档
RED    = RGBColor(0xff,0x5a,0x5f)   # 涨红
GREEN  = RGBColor(0x1e,0xd7,0x60)   # 跌绿
ACCENT = RGBColor(0x00,0xd4,0xff)   # 强调蓝青
TEXT   = RGBColor(0xe0,0xe6,0xf0)   # 主文字
MUTED  = RGBColor(0x8a,0x96,0xb0)   # 次要
LINE   = RGBColor(0x2a,0x35,0x4d)   # 分割线
FONT   = "Microsoft YaHei"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ======== 基础工具函数 ========
def slide():
    """创建空白幻灯片，自动铺满深色背景"""
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0,0,SW,SH)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    return s

def rect(s, x,y,w,h, fill=None, line=LINE, line_w=0.75, rounded=False):
    """绘制矩形块"""
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x,y,w,h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def txt(s, x,y,w,h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """
    添加文本框，runs 格式：
    runs = [
        [[(text, fontSize, color, bold), ...]],   # 第1段，可多色
        [[(text, fontSize, color, bold), ...]],   # 第2段
    ]
    单段简化：runs = [(text, fontSize, color, bold)]
    """
    tb = s.shapes.add_textbox(x,y,w,h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left=Inches(0.05); tf.margin_right=Inches(0.05)
    tf.margin_top=Inches(0.02); tf.margin_bottom=Inches(0.02)
    if isinstance(runs, tuple): runs = [runs]
    for i,line in enumerate(runs):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(line, tuple): line = [line]
        for (t,sz,col,bold) in line:
            r = p.add_run(); r.text=t
            r.font.size=Pt(sz); r.font.color.rgb=col
            r.font.bold=bold; r.font.name=FONT
    return tb

def kicker(s, x, y, t, size=16):
    """章节标签"""
    txt(s, x, y, Inches(6), Inches(0.35), [[(t.upper(), size, ACCENT, True)]])

def title(s, x, y, t, size=40):
    """页面标题"""
    txt(s, x, y, Inches(11.5), Inches(0.9), [[(t, size, TEXT, True)]])

def pagenum(s, n, total):
    """页码：n/total"""
    txt(s, SW-Inches(1.4), Inches(0.3), Inches(1.1), Inches(0.35),
        [[(f"{n:02d} / {total:02d}", 14, MUTED, False)]], align=PP_ALIGN.RIGHT)

def card(s, x, y, w, h, num, head, body, accent=ACCENT):
    """
    通用卡片：num(编号) + head(标题) + body(描述)
    字号：num=14pt, head=18pt, body=15pt
    """
    rect(s, x, y, w, h, fill=CARD, line=LINE, line_w=0.75, rounded=True)
    txt(s, x+Inches(0.2), y+Inches(0.14), Inches(4), Inches(0.32),
        [[(num, 14, accent, True)]])
    txt(s, x+Inches(0.2), y+Inches(0.5), w-Inches(0.4), Inches(0.44),
        [[(head, 18, TEXT, True)]])
    txt(s, x+Inches(0.2), y+Inches(1.0), w-Inches(0.4), h-Inches(1.1),
        [[(body, 15, MUTED, False)]])

# ======== 卡片网格布局辅助 ========
def card_grid(s, items, cols, x=Inches(0.9), y=Inches(2.0),
              cw=Inches(3.78), ch=Inches(1.6), gpx=Inches(0.18), gpy=Inches(0.18),
              accent=ACCENT):
    """
    items = [(num, head, body), ...], cols=3/4
    自动计算行列位置，绘制卡片网格
    """
    for i,(num,head,body) in enumerate(items):
        r,c = divmod(i, cols)
        card(s, x+c*(cw+gpx), y+r*(ch+gpy), cw, ch, num, head, body, accent)

# ======== 表格辅助 ========
def build_table(s, headers, rows, x, y, tw, colw, rh=Inches(0.6)):
    """
    headers = [col1, col2, ...]
    rows = [[val1, val2, ...], ...]
    colw = [Inches(w1), Inches(w2), ...]  # 必须是整数 Inches
    rh = 行高
    返回 table 对象供进一步设置
    """
    tbl = s.shapes.add_table(len(rows)+1, len(headers), x, y, tw, rh*(len(rows)+1)).table
    for ci,w in enumerate(colw):
        tbl.columns[ci].width = w
    # 表头
    for ci,h in enumerate(headers):
        c = tbl.cell(0,ci); c.fill.solid(); c.fill.fore_color.rgb = CARD2
        c.margin_left=Inches(0.12); c.vertical_anchor=MSO_ANCHOR.MIDDLE
        p=c.text_frame.paragraphs[0]; r=p.add_run(); r.text=h
        r.font.name=FONT; r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=ACCENT
    # 数据行
    for ri,row in enumerate(rows):
        tbl.rows[ri+1].height = rh
        for ci,val in enumerate(row):
            c = tbl.cell(ri+1,ci); c.fill.solid(); c.fill.fore_color.rgb = CARD
            c.margin_left=Inches(0.12); c.vertical_anchor=MSO_ANCHOR.MIDDLE
            p=c.text_frame.paragraphs[0]; r=p.add_run(); r.text=str(val)
            r.font.name=FONT; r.font.size=Pt(15.5); r.font.color.rgb=TEXT
    return tbl

# ================================================================
# ======== 下方是幻灯片内容区 —— 逐页替换你的内容 ========
# ================================================================

TOTAL = 3  # 总页数，用于页码

# ---- SLIDE 1 封面 ----
s = slide()
kicker(s, Inches(0.9), Inches(2.0), "YOUR KICKER")
txt(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.6),
    [[("报告主标题", 52, TEXT, True)]])
txt(s, Inches(0.9), Inches(4.2), Inches(11), Inches(1.2),
    [[("副标题 / 描述信息", 22, MUTED, False)],
     [("数据源 · 日期", 20, MUTED, False)]])
pagenum(s, 1, TOTAL)

# ---- SLIDE 2 卡片示例 ----
s = slide()
kicker(s, Inches(0.9), Inches(0.7), "Section Label")
title(s, Inches(0.9), Inches(1.1), "页面标题")
card_grid(s, [
    ("01", "卡片标题 A", "这里写卡片的内容描述"),
    ("02", "卡片标题 B", "这里写卡片的内容描述"),
    ("03", "卡片标题 C", "这里写卡片的内容描述"),
], cols=3)
pagenum(s, 2, TOTAL)

# ---- SLIDE 3 表格示例 ----
s = slide()
kicker(s, Inches(0.9), Inches(0.7), "Data Table")
title(s, Inches(0.9), Inches(1.1), "表格页面")
build_table(s,
    ["列1", "列2"],
    [["数据A1","数据A2"], ["数据B1","数据B2"], ["数据C1","数据C2"]],
    x=Inches(0.9), y=Inches(2.0), tw=Inches(11.5),
    colw=[Inches(4), Inches(7)],
    rh=Inches(0.55))
txt(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.5),
    [[("免责声明：本报告仅供学习研究，不构成投资建议。", 16, MUTED, True)]])
pagenum(s, 3, TOTAL)

# ================================================================
# ======== 输出 ========
# ================================================================
out = r"C:\Users\Administrator\WorkBuddy\Claw\output.pptx"
prs.save(out)
print(f"Saved: {out}  |  Slides: {len(prs.slides._sldIdLst)}")
